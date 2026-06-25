from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

INPUT_CSV = "ci_test_priority_map_10_14_0_0_ci_2026-05-20_history_from_10_9_0_0_tuned_new_tests.csv"
INPUT_CHART = "priority_p4_scenarios.png"
OUTPUT_PPT = "priority_strategy_manager_briefing.pptx"


def build_scenarios(df):
    df = df.copy()
    for c in [
        "hist_failed",
        "hist_executions",
        "failed",
        "executions",
        "avg_runtime_minutes",
        "priority_score",
        "hw_accel_relevance",
    ]:
        df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)

    df["hist_fail_rate"] = (df["hist_failed"] / df["hist_executions"].replace(0, pd.NA)).fillna(0)
    df["curr_fail_rate"] = (df["failed"] / df["executions"].replace(0, pd.NA)).fillna(0)
    df["cell"] = df["platform_type"].astype(str) + " | " + df["mode"].astype(str)

    runners_per_cell = 4
    base_wall_h = (df.groupby("cell")["avg_runtime_minutes"].sum() / runners_per_cell).max() / 60
    total_tests = len(df)
    total_runtime_h = df["avg_runtime_minutes"].sum() / 60

    scenarios = {
        "Strict": (df["priority_score"] < 22)
        & (df["hist_fail_rate"] < 0.01)
        & (df["curr_fail_rate"] == 0)
        & (df["hw_accel_relevance"] == 0),
        "Balanced": (df["priority_score"] < 25)
        & (df["hist_fail_rate"] < 0.02)
        & (df["curr_fail_rate"] < 0.01)
        & (df["hw_accel_relevance"] == 0),
        "Aggressive": (df["priority_score"] < 28)
        & (df["hist_fail_rate"] < 0.03)
        & (df["curr_fail_rate"] < 0.02)
        & (df["hw_accel_relevance"] == 0),
    }

    rows = []
    for name, mask in scenarios.items():
        p4 = df[mask]
        kept = df[~mask]
        wall_h = (kept.groupby("cell")["avg_runtime_minutes"].sum() / runners_per_cell).max() / 60

        hist_failed_exec_in_p4 = int(p4["hist_failed"].sum())
        curr_failed_exec_in_p4 = int(p4["failed"].sum())

        if name == "Strict":
            risk = "Low"
            rationale = "Small defer set; near-zero current failures and negligible historical failure volume."
        elif name == "Balanced":
            risk = "Medium-Low"
            rationale = "Meaningful savings with controlled historical failure exposure."
        else:
            risk = "Medium"
            rationale = "Highest savings but larger deferred historical-failure footprint."

        rows.append(
            {
                "scenario": name,
                "tests_p4": len(p4),
                "tests_p4_pct": len(p4) / total_tests * 100,
                "runtime_saved_h": p4["avg_runtime_minutes"].sum() / 60,
                "runtime_saved_days": p4["avg_runtime_minutes"].sum() / 60 / 24,
                "wall_saved_h": base_wall_h - wall_h,
                "wall_saved_days": (base_wall_h - wall_h) / 24,
                "hist_failed_exec_in_p4": hist_failed_exec_in_p4,
                "curr_failed_exec_in_p4": curr_failed_exec_in_p4,
                "risk": risk,
                "risk_rationale": rationale,
            }
        )

    out = pd.DataFrame(rows)
    return out


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.4), Inches(0.6))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 55, 99)


def add_subtitle(slide, text):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(12.4), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(80, 80, 80)


def add_bullet_slide(prs, title, subtitle, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title)
    add_subtitle(slide, subtitle)

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.8), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(14)


def add_table_slide(prs, title, subtitle, scenario_df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title)
    add_subtitle(slide, subtitle)

    rows = 1 + len(scenario_df)
    cols = 8
    table = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.5), Inches(12.5), Inches(3.3)).table

    headers = [
        "Approach",
        "P4 Tests",
        "P4 %",
        "Runtime Saved (days)",
        "Wall Saved (days)",
        "Current Fails in P4",
        "Historical Fails in P4",
        "Risk",
    ]

    widths = [1.5, 1.2, 1.0, 2.0, 1.8, 1.6, 1.8, 1.2]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(31, 55, 99)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    for r, (_, row) in enumerate(scenario_df.iterrows(), start=1):
        vals = [
            row["scenario"],
            f"{int(row['tests_p4']):,}",
            f"{row['tests_p4_pct']:.1f}%",
            f"{row['runtime_saved_days']:.1f}",
            f"{row['wall_saved_days']:.2f}",
            f"{int(row['curr_failed_exec_in_p4'])}",
            f"{int(row['hist_failed_exec_in_p4']):,}",
            row["risk"],
        ]
        for c, v in enumerate(vals):
            cell = table.cell(r, c)
            cell.text = str(v)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 248, 252)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.CENTER

    note = slide.shapes.add_textbox(Inches(0.4), Inches(5.1), Inches(12.4), Inches(1.6))
    nt = note.text_frame
    np = nt.paragraphs[0]
    np.text = "Assumptions: 4 runners per platform+mode cell (6 cells total), 2-week sprint running 24/7, baseline executes P0+P1+P2+P3, P4 is deferred per approach criteria."
    np.font.size = Pt(12)
    np.font.color.rgb = RGBColor(70, 70, 70)


def add_chart_slide(prs, title, subtitle, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title)
    add_subtitle(slide, subtitle)
    slide.shapes.add_picture(image_path, Inches(0.35), Inches(1.4), Inches(12.6), Inches(5.7))


def add_calculation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Priority System Calculation")
    add_subtitle(slide, "Weighted score model and priority band mapping")

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(2.3))
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Score = 100 * (3*ChangeImpact + 3*HistoricalFailure + 2*PlatformModeRisk + 1*HWAccelRelevance + 1*BusinessCriticality - 1*RuntimeCost) / 9"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 55, 99)

    details = [
        "HistoricalFailure blends long-term fail rate, recent 14-day fail rate, and historical fail volume.",
        "PlatformModeRisk defaults by platform and mode (FPGA/Software/EZchip x Routing/Transparent).",
        "RuntimeCost is normalized so long-running tests are slightly de-prioritized.",
        "Band thresholds: P0-Critical > 79.99, P1-High 55.00-79.99, P2-Medium 40.00-54.99, P3-Low <= 39.99.",
    ]

    dbox = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(11.9), Inches(2.6))
    dtf = dbox.text_frame
    dtf.word_wrap = True
    for i, line in enumerate(details):
        dp = dtf.paragraphs[0] if i == 0 else dtf.add_paragraph()
        dp.text = f"- {line}"
        dp.font.size = Pt(15)
        dp.font.color.rgb = RGBColor(45, 45, 45)
        dp.space_after = Pt(8)


def add_new_test_policy_slide(prs, df):
    total = len(df)

    is_new = pd.Series([False] * total)
    stable = pd.Series([False] * total)
    promoted = pd.Series([False] * total)

    if "is_new_test" in df.columns:
        is_new = df["is_new_test"].astype(str).str.lower().eq("true")
    if "new_test_stable" in df.columns:
        stable = df["new_test_stable"].astype(str).str.lower().eq("true")
    if "new_test_policy_applied" in df.columns:
        promoted = df["new_test_policy_applied"].astype(str).str.lower().eq("true")

    new_count = int(is_new.sum())
    stable_count = int(stable.sum())
    promoted_count = int(promoted.sum())

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "New Test Policy")
    add_subtitle(slide, "New tests are promoted to P1 until they demonstrate stability")

    bullets = [
        "Policy rule: if a test is new and not stable, force minimum band to P1-High.",
        "Stability rule: stable only when executions >= 3 and failures == 0 in current window.",
        f"Current run totals: {new_count:,} new tests out of {total:,} rows.",
        f"Policy applied promotions: {promoted_count:,} rows promoted to P1-High.",
        f"Stable new tests in current window: {stable_count:,} rows.",
    ]

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.8), Inches(5.4))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"- {line}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(12)


def add_risk_slide(prs, title, subtitle, scenario_df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title)
    add_subtitle(slide, subtitle)

    y = 1.6
    for _, row in scenario_df.iterrows():
        box = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(12.1), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 248, 252)
        box.line.color.rgb = RGBColor(210, 220, 235)

        tx = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.15), Inches(11.5), Inches(1.2))
        tf = tx.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = f"{row['scenario']} — Risk: {row['risk']}"
        p1.font.bold = True
        p1.font.size = Pt(16)
        p1.font.color.rgb = RGBColor(31, 55, 99)

        p2 = tf.add_paragraph()
        p2.text = (
            f"Deferred tests: {int(row['tests_p4']):,} ({row['tests_p4_pct']:.1f}%), "
            f"wall-clock saved: {row['wall_saved_days']:.2f} days, "
            f"historical fails deferred: {int(row['hist_failed_exec_in_p4']):,}, "
            f"current fails deferred: {int(row['curr_failed_exec_in_p4'])}."
        )
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(60, 60, 60)

        p3 = tf.add_paragraph()
        p3.text = row["risk_rationale"]
        p3.font.size = Pt(12)
        p3.font.color.rgb = RGBColor(80, 80, 80)

        y += 1.75


def add_recommendation_slide(prs, title, subtitle):
    bullets = [
        "Run full P0 + P1 + P2 + P3 each sprint (fits with 4 runners per cell under 24/7 operation).",
        "Apply new-test P1 policy consistently: new tests stay at P1 until they are stable.",
        "Use P4 defer scenarios as optional optimization, not as a requirement to fit sprint capacity.",
        "Balanced remains the default optimization mode when extra buffer is needed.",
        "Promote any deferred test to P2 immediately on first failure.",
        "Run a full P3/P4 sweep before release milestones to prevent blind spots.",
        "Track escaped defects and rebalance thresholds every 2-3 sprints.",
    ]
    add_bullet_slide(prs, title, subtitle, bullets)


def main():
    df = pd.read_csv(INPUT_CSV)
    scenario_df = build_scenarios(df)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_bullet_slide(
        prs,
        "Automation Regression Optimization Plan",
        f"Prepared for QA Management | {datetime.now().strftime('%Y-%m-%d')}",
        [
            "Goal: reduce regression cycle duration while preserving release quality.",
            "Baseline run policy: execute P0 + P1 + P2 + P3 regularly.",
            "New-test policy: new tests are treated as P1 until they meet stability criteria.",
            "Capacity assumption: 4 runners per platform+mode cell, 24/7 during 2-week sprint.",
            "Optimization lever: defer ultra-low-risk tests into optional P4 buckets when extra buffer is required.",
            "Scenarios evaluated: Strict, Balanced, Aggressive.",
        ],
    )

    add_calculation_slide(prs)

    add_new_test_policy_slide(prs, df)

    add_table_slide(
        prs,
        "Scenario Comparison (Strict vs Balanced vs Aggressive)",
        "Time reduction and risk exposure when introducing optional P4",
        scenario_df,
    )

    add_chart_slide(
        prs,
        "Visual Comparison",
        "P4 coverage, runtime saved, and wall-clock saved by approach",
        INPUT_CHART,
    )

    add_risk_slide(
        prs,
        "Risk Analysis",
        "What could be missed when deferred tests are not executed in each cycle",
        scenario_df,
    )

    add_recommendation_slide(
        prs,
        "Recommended Rollout",
        "Pragmatic path to reduce cycle time while controlling quality risk",
    )

    prs.save(OUTPUT_PPT)
    scenario_df.to_csv("priority_strategy_ppt_data.csv", index=False)
    print(f"Created {OUTPUT_PPT}")
    print("Created priority_strategy_ppt_data.csv")


if __name__ == "__main__":
    main()
