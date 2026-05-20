"""
PowerPoint Report Generator for DefensePro Weekly Report.
Embeds Plotly chart images and key data into a multi-slide presentation.
"""
import io
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# Slide dimensions (widescreen 10x7.5)
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

if HAS_PPTX:
    BLUE = RGBColor(0x19, 0x76, 0xD2)
    DARK = RGBColor(0x33, 0x33, 0x33)
    GRAY = RGBColor(0x66, 0x66, 0x66)
    LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
    GREEN = RGBColor(0x27, 0xAE, 0x60)
    RED = RGBColor(0xE7, 0x4C, 0x3C)
    ORANGE = RGBColor(0xF5, 0x76, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _title_bar(slide, title):
    """Add a blue title bar at the top of a slide."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE


def _add_image_slide(prs, title, img_bytes, subtitle=None):
    """Add a slide with a full-width chart image."""
    if not img_bytes:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, title)

    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(9), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY

    img_top = Inches(1.5) if subtitle else Inches(1.1)
    img_height = Inches(5.8) if subtitle else Inches(6.2)
    stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(stream, Inches(0.3), img_top, Inches(9.4), img_height)


def _slide_title(prs, version, sprint_name, period, ci_run_start):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(3.2))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(8.8), Inches(2.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"DefensePro {version}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "Weekly Status Report"
    p2.font.size = Pt(22)
    p2.font.color.rgb = WHITE

    txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.6), Inches(8.8), Inches(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    lines = [
        f"Sprint: {sprint_name}",
        f"Period: {period}",
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
    ]
    if ci_run_start:
        lines.append(f"CI Cycle Start: {ci_run_start}")
    for i, line in enumerate(lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)


def _slide_summary(prs, data):
    """Executive summary with key metric cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Executive Summary")

    metrics = [
        ("Bugs on Dev", str(data['bugs_on_dev']), ORANGE if data['bugs_on_dev'] > 0 else GREEN),
        ("Bugs on QA", str(data['bugs_on_qa']), ORANGE if data['bugs_on_qa'] > 3 else BLUE),
        ("Closed This Week", str(data['bugs_closed_week']), GREEN),
        ("Tests Executed", f"{data['total_tests']:,}", BLUE),
        ("Pass Ratio", f"{data['pass_ratio']:.1f}%", GREEN if data['pass_ratio'] >= 90 else RED),
        ("Coverage", f"{data.get('coverage', 0):.1f}%", BLUE),
    ]

    for i, (label, value, color) in enumerate(metrics):
        row = i // 3
        col = i % 3
        left = Inches(0.4 + col * 3.15)
        top = Inches(1.3 + row * 2.8)

        card = slide.shapes.add_shape(5, left, top, Inches(2.9), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

        txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.4), Inches(2.5), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        txBox2 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.5), Inches(2.5), Inches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER


def _slide_bugs_table(prs, data):
    """Bugs on Dev + QA table."""
    dev_list = data.get('bugs_on_dev_list', [])
    qa_list = data.get('bugs_on_qa_list', [])
    all_bugs = []
    for b in dev_list[:7]:
        all_bugs.append(('Dev', b['key'], b['priority'], b['summary'][:55]))
    for b in qa_list[:7]:
        all_bugs.append(('QA', b['key'], b['priority'], b['summary'][:55]))

    if not all_bugs:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, f"Open Bugs — {data['bugs_on_dev']} Dev / {data['bugs_on_qa']} QA")

    rows = len(all_bugs) + 1
    tbl = slide.shapes.add_table(rows, 4, Inches(0.3), Inches(1.2), Inches(9.4), Inches(0.45 * rows)).table
    tbl.columns[0].width = Inches(0.7)
    tbl.columns[1].width = Inches(1.3)
    tbl.columns[2].width = Inches(1.2)
    tbl.columns[3].width = Inches(6.2)

    headers = ['Phase', 'Key', 'Priority', 'Summary']
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for par in cell.text_frame.paragraphs:
            par.font.color.rgb = WHITE
            par.font.size = Pt(10)
            par.font.bold = True

    for ri, (phase, key, prio, summary) in enumerate(all_bugs, 1):
        tbl.cell(ri, 0).text = phase
        tbl.cell(ri, 1).text = key
        tbl.cell(ri, 2).text = prio
        tbl.cell(ri, 3).text = summary
        for ci in range(4):
            for par in tbl.cell(ri, ci).text_frame.paragraphs:
                par.font.size = Pt(9)
        prio_color = RED if prio in ('High', 'Highest', 'Critical') else ORANGE if prio == 'Medium' else DARK
        for par in tbl.cell(ri, 2).text_frame.paragraphs:
            par.font.color.rgb = prio_color
            par.font.bold = True


def _slide_automation_summary(prs, data):
    """CI / Automation key stats."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "CI Iteration — Automation Status")

    lines = [
        f"Tests Executed: {data['total_tests']:,} unique tests, {data['total_executions']:,} total runs",
        f"Pass Ratio: {data['pass_ratio']:.1f}% ({data['passed']:,} passed / {data['failed']:,} failed)",
        f"Coverage: {data.get('coverage', 0):.1f}%",
    ]
    if data.get('critical_failures', 0) > 0:
        lines.append(f"Critical Failures: {data['critical_failures']} tests failing on ALL platforms")
    if data.get('avg_daily_rate'):
        lines.append(f"Velocity: {data['avg_daily_rate']}%/day")
    if data.get('projected_coverage'):
        lines.append(f"Projected at Sprint End: {data['projected_coverage']}%")

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {line}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.space_after = Pt(14)


def _slide_sub_exec(prs, data):
    """Sub test execution progress."""
    total = data.get('sub_exec_total', 0)
    if total == 0:
        return
    completed = data.get('sub_exec_completed', 0)
    in_progress = data.get('sub_exec_in_progress', 0)
    not_started = data.get('sub_exec_not_started', 0)
    xray = data.get('xray_summary', {})

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Sub Test Execution Status")

    bar_left = Inches(0.5)
    bar_top = Inches(1.5)
    bar_w = Inches(9)
    bar_h = Inches(0.6)

    bg_bar = slide.shapes.add_shape(5, bar_left, bar_top, bar_w, bar_h)
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    bg_bar.line.fill.background()

    pct = completed / total if total > 0 else 0
    if pct > 0:
        fg_bar = slide.shapes.add_shape(5, bar_left, bar_top, Inches(9 * pct), bar_h)
        fg_bar.fill.solid()
        fg_bar.fill.fore_color.rgb = GREEN
        fg_bar.line.fill.background()

    txBar = slide.shapes.add_textbox(bar_left, bar_top + Inches(0.05), bar_w, bar_h)
    tf = txBar.text_frame
    p = tf.paragraphs[0]
    p.text = f"{pct*100:.0f}% Complete ({completed}/{total})"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE if pct > 0.4 else DARK
    p.alignment = PP_ALIGN.CENTER

    lines = [
        f"Completed: {completed}  |  In Progress: {in_progress}  |  Not Started: {not_started}",
    ]
    if xray.get('total_tests', 0) > 0:
        lines.append(f"Xray Tests: {xray['total_tests']}  |  Coverage: {xray.get('testing_coverage', 0):.1f}%  |  Pass Ratio: {xray.get('pass_ratio', 0):.1f}%")
        lines.append(f"Automation Coverage: {xray.get('automation_coverage', 0):.1f}%")

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.space_after = Pt(12)


def _slide_changelog(prs, build_changelogs):
    """Build changelog."""
    if not build_changelogs:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    total_changes = sum(len(b['changes']) for b in build_changelogs)
    _title_bar(slide, f"Build Changes — {total_changes} product commits across {len(build_changelogs)} builds")

    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(9.2), Inches(6))
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for build in build_changelogs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        icon = "✓" if build.get('result') == 'SUCCESS' else "✗"
        p.text = f"{icon} {build['displayName']} — {len(build['changes'])} change(s)"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = GREEN if build.get('result') == 'SUCCESS' else RED
        p.space_before = Pt(10)

        if not build.get('available', True):
            p = tf.add_paragraph()
            p.text = f"    Jenkins build unavailable (HTTP {build.get('http_status') or 'unknown'})"
            p.font.size = Pt(9)
            p.font.color.rgb = GRAY
            continue

        for change in build['changes'][:6]:
            p = tf.add_paragraph()
            p.text = f"    {change['commitId']}  {change['msg'][:65]}  ({change['author']})"
            p.font.size = Pt(9)
            p.font.color.rgb = GRAY

        if not build['changes']:
            p = tf.add_paragraph()
            skipped = build.get('skipped_change_count', 0)
            suffix = f" ({skipped} CI trigger change(s) hidden)" if skipped else ""
            p.text = f"    No product changes recorded{suffix}"
            p.font.size = Pt(9)
            p.font.color.rgb = GRAY

        if len(build['changes']) > 6:
            p = tf.add_paragraph()
            p.text = f"    ... and {len(build['changes']) - 6} more"
            p.font.size = Pt(9)
            p.font.color.rgb = GRAY


def _slide_insights(prs, insights, ai_insights):
    """Insights slide."""
    if not insights and not ai_insights:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Key Insights")

    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(9.2), Inches(6))
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    if insights:
        for insight in insights[:8]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = f"•  {insight}"
            p.font.size = Pt(12)
            p.font.color.rgb = DARK
            p.space_after = Pt(6)

    if ai_insights:
        p = tf.add_paragraph()
        p.space_before = Pt(14)
        p = tf.add_paragraph()
        p.text = "AI-Generated Insights:"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.space_after = Pt(6)

        import re
        clean = re.sub(r'<[^>]+>', '', ai_insights)
        for line in clean.split('\n')[:8]:
            line = line.strip()
            if line:
                p = tf.add_paragraph()
                p.text = line[:120]
                p.font.size = Pt(10)
                p.font.color.rgb = GRAY


def _slide_platform_table(prs, platform_type_data):
    """Platform type coverage table."""
    if not platform_type_data:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Platform Type & Mode Coverage")

    sorted_data = sorted(platform_type_data, key=lambda x: x['platform_type_mode'])
    rows = min(len(sorted_data), 12) + 1
    tbl = slide.shapes.add_table(rows, 5, Inches(0.3), Inches(1.2), Inches(9.4), Inches(0.4 * rows)).table
    tbl.columns[0].width = Inches(3.0)
    tbl.columns[1].width = Inches(1.5)
    tbl.columns[2].width = Inches(1.5)
    tbl.columns[3].width = Inches(1.5)
    tbl.columns[4].width = Inches(1.9)

    headers = ['Platform Type & Mode', 'Tests', 'Baseline', 'Coverage', 'Pass Ratio']
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for par in cell.text_frame.paragraphs:
            par.font.color.rgb = WHITE
            par.font.size = Pt(10)
            par.font.bold = True

    for ri, p in enumerate(sorted_data[:12], 1):
        tbl.cell(ri, 0).text = p['platform_type_mode']
        tbl.cell(ri, 1).text = str(p['tests'])
        tbl.cell(ri, 2).text = str(p['available_tests'])
        tbl.cell(ri, 3).text = f"{p['coverage']:.1f}%"
        tbl.cell(ri, 4).text = f"{p['pass_ratio']:.1f}%"
        for ci in range(5):
            for par in tbl.cell(ri, ci).text_frame.paragraphs:
                par.font.size = Pt(9)
        cov_color = RED if p['coverage'] < 70 else ORANGE if p['coverage'] < 90 else GREEN
        for par in tbl.cell(ri, 3).text_frame.paragraphs:
            par.font.color.rgb = cov_color
            par.font.bold = True


def fig_to_png(fig, width=1400, height=700):
    """Convert a Plotly figure to PNG bytes. Returns None on failure."""
    if fig is None:
        return None
    try:
        return fig.to_image(format="png", width=width, height=height, scale=2)
    except Exception:
        return None


def generate_ppt(report_data, output_path):
    """
    Generate a PowerPoint presentation from report data.
    
    Chart images should be passed as PNG bytes under keys:
        chart_automation, chart_bugs, chart_historical, chart_high_sev,
        chart_sub_exec, chart_xray_exec, chart_xray_method
    
    Returns output_path if successful, None if python-pptx unavailable.
    """
    if not HAS_PPTX:
        print("⚠️  python-pptx not installed, skipping PPT generation")
        return None

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Title
    _slide_title(prs, report_data['version'], report_data['sprint_name'],
                 report_data['period'], report_data.get('ci_run_start'))

    # 2. Executive Summary
    _slide_summary(prs, report_data)

    # 3. Automation Status text
    _slide_automation_summary(prs, report_data)

    # 4. Automation chart
    _add_image_slide(prs, "Automation Results by Platform", report_data.get('chart_automation'))

    # 5. Platform type table
    _slide_platform_table(prs, report_data.get('platform_type_data'))

    # 6. Bug Status chart
    _add_image_slide(prs, "Bug Status Distribution", report_data.get('chart_bugs'))

    # 7. Historical Bug Trend
    _add_image_slide(prs, "Historical Bug Trend", report_data.get('chart_historical'))

    # 8. High Severity Trend
    _add_image_slide(prs, "High/Critical Priority Bug Trend", report_data.get('chart_high_sev'))

    # 9. Bugs table
    _slide_bugs_table(prs, report_data)

    # 10. Sub Test Execution
    _slide_sub_exec(prs, report_data)

    # 11. Sub exec chart
    _add_image_slide(prs, "Sub Test Execution Progress", report_data.get('chart_sub_exec'))

    # 12. Xray charts
    _add_image_slide(prs, "Xray Execution Rate", report_data.get('chart_xray_exec'))
    _add_image_slide(prs, "Xray Test Method Distribution", report_data.get('chart_xray_method'))

    # 13. Build Changelog
    _slide_changelog(prs, report_data.get('build_changelogs'))

    # 14. Insights
    _slide_insights(prs, report_data.get('insights'), report_data.get('ai_insights'))

    prs.save(output_path)
    return output_path
